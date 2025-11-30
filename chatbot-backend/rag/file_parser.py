# rag/file_parser.py
import os
import io
import re
import fitz  # PyMuPDF
import pdfplumber
import docx
import pandas as pd
from pptx import Presentation

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".csv", ".xlsx", ".txt"}

def detect_ext(path: str) -> str:
    return os.path.splitext(path)[-1].lower()

def clean_text(t: str) -> str:
    # Basic cleanup
    t = t.replace("\u00a0", " ")
    t = re.sub(r"[ \t]+", " ", t)
    t = re.sub(r"\n{3,}", "\n\n", t)
    return t.strip()

def parse_pdf(path: str) -> str:
    # Try pdfplumber for text + tables
    try:
        all_text = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                txt = page.extract_text() or ""
                all_text.append(txt)
        return clean_text("\n\n".join(all_text))
    except Exception:
        # fallback to PyMuPDF
        doc = fitz.open(path)
        text = "\n\n".join(page.get_text() for page in doc)
        return clean_text(text)

def parse_docx(path: str) -> str:
    d = docx.Document(path)
    parts = []
    for p in d.paragraphs:
        parts.append(p.text)
    # tables as CSV-like lines
    for table in d.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return clean_text("\n".join(parts))

def parse_pptx(path: str) -> str:
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slide_text = f"[Slide {i}]\n" + "\n".join(texts)
        slides.append(slide_text)
    return clean_text("\n\n".join(slides))

def parse_csv(path: str) -> str:
    df = pd.read_csv(path, nrows=2000)  # cap for sanity
    return clean_text(df.to_csv(index=False))

def parse_xlsx(path: str) -> str:
    dfs = pd.read_excel(path, sheet_name=None)
    parts = []
    for sheet_name, df in dfs.items():
        parts.append(f"[Sheet: {sheet_name}]")
        parts.append(df.head(2000).to_csv(index=False))
    return clean_text("\n\n".join(parts))

def parse_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return clean_text(f.read())

def parse_file(path: str) -> str:
    ext = detect_ext(path)
    if ext not in SUPPORTED_EXTS:
        raise ValueError(f"Unsupported file type: {ext}")
    if ext == ".pdf":
        return parse_pdf(path)
    if ext == ".docx":
        return parse_docx(path)
    if ext == ".pptx":
        return parse_pptx(path)
    if ext == ".csv":
        return parse_csv(path)
    if ext == ".xlsx":
        return parse_xlsx(path)
    if ext == ".txt":
        return parse_txt(path)
